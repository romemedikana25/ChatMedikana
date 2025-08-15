# Extract Text Import
from pptx import Presentation
from docx import Document
import pandas as pd
from pypdf import PdfReader
import re
import unicodedata

# Translator Import
from langdetect import detect
from deep_translator import GoogleTranslator

from global_vars import AMERICAS, CODES


#----- Helper Functions for LLM Scripts ------

# ---- TRANSLATION ----
def translate_text(text, target_language='en'):
    try: 
        return GoogleTranslator(source='auto', target=target_language).translate(text)
    except Exception as e:
        print(f"Translation error: {e}")
        return text
    
# ---- TEXT NORMALIZATION ----
def normalize_text(s: str) -> str:
    if not s:
        return s
    s = s.strip().lower()
    # remove surrounding quotes
    s = s.strip("'\"“”‘’")
    # remove diacritics
    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch for ch in s if not unicodedata.combining(ch))
    # collapse spaces
    s = re.sub(r"\s+", " ", s)
    return s

# ---- FILE LOADERS ----
# Extract Text from PDF Files
def extract_pdf_text(filepath):
    try:
        text = ""
        with open(filepath, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return translate_text(text)
    except Exception as e:
        print(f"Error reading PDF {filepath}: {e}")
        return ""

# Extract Text from PowerPoint Presentations
def extract_pptx_text(filepath):
    prs = Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    joined_text = "\n".join(text)
    return translate_text(joined_text)

# Extract Text from Word Documents
def extract_docx_text(filepath):
    doc = Document(filepath)
    text = "\n".join([para.text for para in doc.paragraphs])
    return translate_text(text)

# Excel Text Extraction
def extract_xlsx_text(filepath):
    # Skip temporary lock files created by Excel
    if filepath.name.startswith("~$"):
        return ""

    try:
        dfs = pd.read_excel(filepath, sheet_name=None, engine="openpyxl")
        text_blocks = []
        for sheet, df in dfs.items():
            text_blocks.append(f"--- Sheet: {sheet} ---")
            text_blocks.append(df.to_string(index=False))
        text = "\n".join(text_blocks)
        return translate_text(text)
    except Exception as e:
        print(f"Error reading Excel {filepath}: {e}")
        return ""



# --- Infering Meta Data Helpers ---

def infer_counties_from_file(text):
    countries = []
    for country in AMERICAS:
        if country.lower() in text.lower():
            countries.append(country)
    return countries

def extract_country_from_query(query):
    query = query.lower()
    countries = []
    for country in AMERICAS:
        if country.lower() in query:
            countries.append(country)
    return countries


# ---- Intent Detection for Table Queries ----
# def is_table_or_numerical_query(query):
#     table_keywords = [
#         "spreadsheet", "excel", "csv", "dataframe",
#         "columns", "rows", "numeric", "numerical", "quantitative",
#         "mean", "median", "average", "sum", "count", "total",
#         "filter", "sort", "maximum", "minimum"
#     ]
#     return any(keyword in query.lower() for keyword in table_keywords)


TABLE_HINTS_RE = re.compile(
    r"\b(spreadsheet|excel|csv|sheet|table|dataframe|pivot|columns?|rows?|filter|sort)\b",
    flags=re.IGNORECASE,
)

NUMERIC_OPS_RE = re.compile(
    r"\b(mean|median|average|avg|sum|count|total|max(?:imum)?|min(?:imum)?)\b",
    flags=re.IGNORECASE,
)

AUX_NUMERIC_CUES_RE = re.compile(
    r"\b(per|by|group(?:ed)?|breakdown|distribution|over time)\b|\d",
    flags=re.IGNORECASE,
)

def is_table_or_numerical_query(query: str) -> bool:
    q = query.strip()
    if not q:
        return False
    # 1) Explicit tabular hints → true
    if TABLE_HINTS_RE.search(q):
        return True
    # 2) Numeric ops must be standalone words AND accompanied by a numeric cue
    if NUMERIC_OPS_RE.search(q) and AUX_NUMERIC_CUES_RE.search(q):
        return True
    return False



    
